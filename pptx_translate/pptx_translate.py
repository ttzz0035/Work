import csv
import shutil
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
import os
import yaml
import logging
import sys

# --- ベースディレクトリ設定（EXEと外部config運用対応） ---
# PyInstallerで--onefileにすると、EXE実行中は実行箇所（sys.executable）をベースに動作する
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_DIR = os.path.join(BASE_DIR, "config")
LOG_DIR = os.path.join(BASE_DIR, "logs")
USER_PATHS_FILE = os.path.join(BASE_DIR, "user_paths.yaml")
os.makedirs(LOG_DIR, exist_ok=True)

# ログ設定
logger = logging.getLogger("pptx_translate")
logger.setLevel(logging.INFO)

formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")

file_handler = logging.FileHandler(os.path.join(LOG_DIR, "app.log"), encoding="utf-8")
file_handler.setFormatter(formatter)

stream_handler = logging.StreamHandler()
stream_handler.setFormatter(formatter)

logger.addHandler(file_handler)
logger.addHandler(stream_handler)

# --- ファイル操作ユーティリティ ---

def load_yaml(file_path, default_data=None):
    if not os.path.exists(file_path):
        if default_data is not None:
            save_yaml(file_path, default_data)
            return default_data
        else:
            return {}
    with open(file_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)

def save_yaml(file_path, data):
    with open(file_path, "w", encoding="utf-8") as f:
        yaml.dump(data, f, allow_unicode=True)

def check_config_folder(labels):
    config_path = os.path.join(BASE_DIR, "config")
    if not os.path.exists(config_path):
        messagebox.showerror(labels["error_config_folder_not_found_title"], labels["error_config_folder_not_found_message"])
        logger.error(labels["log_config_folder_not_found"])
        raise FileNotFoundError(labels["error_config_folder_not_found_message"])

# --- CSVエクスポート・インポート処理 ---

def collect_texts_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    rows = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text") and shape.text.strip():
                rows.append([
                    slide_idx,
                    shape_idx,
                    shape.text.strip(),
                    ""  # text_english
                ])
    return rows

def export_texts_to_csv(input_path, output_path, labels):
    try:
        logger.info(f"{labels['log_csv_export_start']}: {input_path} → {output_path}")
        rows = collect_texts_from_pptx(input_path)
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f, quoting=csv.QUOTE_ALL)
            writer.writerow(["slide_idx", "shape_idx", "text_japanese", "text_english"])
            writer.writerows(rows)
        logger.info(f"{labels['log_csv_export_success']}: {output_path}")
        messagebox.showinfo(labels["message_title_success"], f"{labels['message_export_success']}:\n{output_path}")
    except Exception as e:
        logging.error(f"{labels['log_csv_export_failure']}: {e}")
        messagebox.showerror(labels["message_title_error"], f"{labels['message_export_error']}:\n{str(e)}")

def load_translations(csv_path):
    translations = {}
    with open(csv_path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            key = (int(row["slide_idx"]), int(row["shape_idx"]))
            translations[key] = row["text_english"]
    return translations

def apply_translations_to_pptx(prs, translations):
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            key = (slide_idx, shape_idx)
            if key in translations and translations[key] and hasattr(shape, "text"):
                shape.text = translations[key]

def import_translated_texts_and_generate_pptx(original_pptx_path, translated_csv_path, output_pptx_path, labels):
    try:
        logger.info(f"{labels['log_pptx_create_start']}: {original_pptx_path} + {translated_csv_path} → {output_pptx_path}")
        shutil.copy2(original_pptx_path, output_pptx_path)
        prs = Presentation(output_pptx_path)
        translations = load_translations(translated_csv_path)
        apply_translations_to_pptx(prs, translations)
        prs.save(output_pptx_path)
        logger.info(f"{labels['log_pptx_create_success']}: {output_pptx_path}")
        messagebox.showinfo(labels["message_title_success"], f"{labels['message_import_success']}:\n{output_pptx_path}")
    except Exception as e:
        logging.error(f"{labels['log_pptx_create_failure']}: {e}")
        messagebox.showerror(labels["message_title_error"], f"{labels['message_import_error']}:\n{str(e)}")

# --- ファイル選択処理 ---

def get_initial_dir(entry_widget, app_settings):
    current_path = entry_widget.get().strip()
    if current_path and os.path.exists(os.path.dirname(current_path)):
        return os.path.dirname(current_path)
    else:
        return app_settings["app"]["default_dir"]

def browse_file(entry_widget, filetypes, key, user_paths, app_settings, labels, save=False):
    try:
        initial_dir = get_initial_dir(entry_widget, app_settings)
        if save:
            path = filedialog.asksaveasfilename(initialdir=initial_dir, defaultextension=filetypes[0][1], filetypes=filetypes)
        else:
            path = filedialog.askopenfilename(initialdir=initial_dir, filetypes=filetypes)
        if path:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, path)
            user_paths[key] = path
            save_yaml(USER_PATHS_FILE, user_paths)
            logger.info(f"{labels['log_file_selected']}: {path}")
    except Exception as e:
        logging.error(f"{labels['log_file_select_error']}: {e}")

# --- UI部品生成 ---

def create_ui_entry(frame, label_text, comps_key, user_path_key, filetypes, comps, user_paths, app_settings, labels, save=False):
    row = comps[comps_key]["row"]
    width = comps[comps_key]["width"]
    tk.Label(frame, text=labels[label_text]).grid(row=row, column=0, sticky="w")
    entry = tk.Entry(frame, width=width)
    entry.grid(row=row, column=1)
    entry.insert(0, user_paths.get(user_path_key, ""))
    tk.Button(frame, text="...", command=lambda: browse_file(entry, filetypes, user_path_key, user_paths, app_settings, labels, save)).grid(row=row, column=2)
    return entry

# --- イベントリスナー ---

def on_mouse_release(event):
    widget = event.widget
    widget_class = widget.winfo_class()
    logging.debug(
        f"マウスリリース: コンポーネント={widget_class}, "
        f"座標=({widget.winfo_x()},{widget.winfo_y()}), "
        f"サイズ=({widget.winfo_width()}x{widget.winfo_height()})"
    )

# --- メイン関数 ---

def main():
    app_settings = load_yaml(os.path.join(CONFIG_DIR, "app_settings.yaml"))
    labels_all = load_yaml(os.path.join(CONFIG_DIR, "labels.yaml"))
    lang = "ja"  # 言語設定
    labels = labels_all[lang]
    user_paths = load_yaml(USER_PATHS_FILE, default_data={
        "input_pptx": "",
        "output_csv": "",
        "original_pptx": "",
        "translated_csv": "",
        "output_pptx": ""
    })

    logger.info(labels["log_startup"])

    root = tk.Tk()
    root.title(labels["app_title"])
    root.geometry(app_settings["app"]["window_size"])

    comps = app_settings["components"]

    frame_export = tk.LabelFrame(root, text=labels["section_export"], padx=10, pady=10)
    frame_export.pack(fill="both", padx=10, pady=10)

    input_pptx_entry = create_ui_entry(frame_export, "label_input_pptx", "input_pptx", "input_pptx", [("PowerPoint Files", "*.pptx")], comps, user_paths, app_settings, labels)
    output_csv_entry = create_ui_entry(frame_export, "label_output_csv", "output_csv", "output_csv", [("CSV Files", "*.csv")], comps, user_paths, app_settings, labels, save=True)

    tk.Button(frame_export, text=labels["button_export"], height=comps["export_button"]["height"],
              command=lambda: export_texts_to_csv(input_pptx_entry.get(), output_csv_entry.get(), labels)
              ).grid(row=comps["export_button"]["row"], column=comps["export_button"]["column"], pady=10)

    frame_import = tk.LabelFrame(root, text=labels["section_import"], padx=10, pady=10)
    frame_import.pack(fill="both", padx=10, pady=10)

    original_pptx_entry = create_ui_entry(frame_import, "label_original_pptx", "original_pptx", "original_pptx", [("PowerPoint Files", "*.pptx")], comps, user_paths, app_settings, labels)
    translated_csv_entry = create_ui_entry(frame_import, "label_translated_csv", "translated_csv", "translated_csv", [("CSV Files", "*.csv")], comps, user_paths, app_settings, labels)
    output_pptx_entry = create_ui_entry(frame_import, "label_output_pptx", "output_pptx", "output_pptx", [("PowerPoint Files", "*.pptx")], comps, user_paths, app_settings, labels, save=True)

    tk.Button(frame_import, text=labels["button_import"], height=comps["import_button"]["height"],
              command=lambda: import_translated_texts_and_generate_pptx(
                  original_pptx_entry.get(), translated_csv_entry.get(), output_pptx_entry.get(), labels
              )).grid(row=comps["import_button"]["row"], column=comps["import_button"]["column"], pady=10)

    root.bind("<ButtonRelease-1>", on_mouse_release)
    root.mainloop()
    logger.info(labels["log_shutdown"])

if __name__ == "__main__":
    main()
